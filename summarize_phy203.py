from pptx import Presentation
import re

# Read PowerPoint
prs = Presentation('PHY 203 LECTURE NOTE II FINAL.pptx')

content = []
for i, slide in enumerate(prs.slides, 1):
    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                slide_text.append(text)
    
    if slide_text:
        content.append({
            'slide': i,
            'text': '\n'.join(slide_text)
        })

# Generate HTML summary
html = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PHY 203 - Lecture Notes Summary</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            color: #333;
        }
        .container {
            max-width: 1000px;
            margin: 0 auto;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
            overflow: hidden;
        }
        .header {
            background: #0a0a23;
            color: white;
            padding: 30px;
            text-align: center;
        }
        .header h1 {
            font-size: 2.5rem;
            margin-bottom: 10px;
        }
        .header p {
            font-size: 1.1rem;
            opacity: 0.9;
        }
        .content {
            padding: 30px;
        }
        .slide {
            background: #f8f9fa;
            border-left: 4px solid #667eea;
            padding: 20px;
            margin-bottom: 20px;
            border-radius: 5px;
        }
        .slide-number {
            color: #667eea;
            font-weight: bold;
            font-size: 0.9rem;
            margin-bottom: 10px;
        }
        .slide-text {
            white-space: pre-wrap;
            line-height: 1.8;
        }
        .footer {
            background: #0a0a23;
            color: white;
            text-align: center;
            padding: 20px;
            font-size: 0.9rem;
        }
        @media print {
            body {
                background: white;
                padding: 0;
            }
            .container {
                box-shadow: none;
            }
        }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>PHY 203</h1>
            <p>Lecture Notes Summary</p>
        </div>
        <div class="content">
"""

for item in content:
    html += f"""
            <div class="slide">
                <div class="slide-number">Slide {item['slide']}</div>
                <div class="slide-text">{item['text']}</div>
            </div>
"""

html += """
        </div>
        <div class="footer">
            PHY 203 Lecture Notes - Generated Summary
        </div>
    </div>
</body>
</html>
"""

# Save HTML file
with open('PHY_203_Summary.html', 'w', encoding='utf-8') as f:
    f.write(html)

print(f"Summary created successfully!")
print(f"Total slides processed: {len(content)}")
print(f"Output file: PHY_203_Summary.html")
